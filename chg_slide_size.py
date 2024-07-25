import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE

def scale_position(position, scale_width, scale_height):
    return (int(position[0] * scale_width), int(position[1] * scale_height))

def scale_size(size, scale_factor):
    return tuple(int(dim * scale_factor) for dim in size)

def scale_shape(shape, scale_width, scale_height):
    shape.left, shape.top = scale_position((shape.left, shape.top), scale_width, scale_height)

    shape.width = int(shape.width * scale_width)
    shape.height = int(shape.height * scale_height)
    
    if shape.has_text_frame:
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    run.font.size = Pt(run.font.size.pt * scale_height)
                run.font.color.rgb = run.font.color.rgb
            
            if paragraph.line_spacing:
                paragraph.line_spacing = Pt(paragraph.line_spacing.pt * scale_height)
            
            paragraph.alignment = paragraph.alignment

def scale_slide_contents(slide, scale_width, scale_height):
    for shape in slide.shapes:
        scale_shape(shape, scale_width, scale_height)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    cell.text_frame.auto_size = MSO_AUTO_SIZE.NONE
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size:
                                run.font.size = Pt(run.font.size.pt * scale_height)
                            run.font.color.rgb = run.font.color.rgb
                        
                        if paragraph.line_spacing:
                            paragraph.line_spacing = Pt(paragraph.line_spacing.pt * scale_height)
                        
                        paragraph.alignment = paragraph.alignment

def resize_pptx(input_file, output_file):
    prs = Presentation(input_file)

    original_width = prs.slide_width
    original_height = prs.slide_height

    new_width = original_width * 3
    new_height = original_height * 3

    scale_width = new_width / original_width
    scale_height = new_height / original_height

    prs.slide_width = new_width
    prs.slide_height = new_height

    for slide in prs.slides:
        scale_slide_contents(slide, scale_width, scale_height)

    prs.save(output_file)

def process_all_pptx_in_folder(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith('.pptx'):
            input_file = os.path.join(folder_path, filename)
            output_file = os.path.join(folder_path, f'resized_{filename}')
            resize_pptx(input_file, output_file)
            print(f'変換完了: {filename} -> resized_{filename}')

# Slideフォルダ内のすべての.pptxファイルを処理
process_all_pptx_in_folder('Slide')
